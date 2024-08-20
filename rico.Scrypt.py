#Bibliotecas
from pptx import Presentation
from pptx.util import Inches
from docx import Document
import pandas as pd
from datetime import datetime
import comtypes.client

#Funcoes
def extract_doc_contents(doc_path):
    document = Document(doc_path)
    paragraphs = []
    tables = []
    
    # Extrair parágrafos
    for paragraph in document.paragraphs:
        parts = paragraph.text.split('\n\n')
        for part in parts:
            clean_part = part.strip()
            if clean_part:
                paragraphs.append(clean_part)
                
    # Extrair tabelas
    for table in document.tables:
        table_data = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(row_data)
        df = pd.DataFrame(table_data)
        tables.append(df)

    return paragraphs, tables

# Função PowerPoint
def substitute_placeholders_ppt(template_path, output_path, replacements, tables):
    presentation = Presentation(template_path)
    
    for i, slide in enumerate(presentation.slides):
        if i >= len(replacements):
            break
        slide_replacements = replacements[i]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for key, value in slide_replacements.items():
                    if key in shape.text:
                        shape.text = shape.text.replace(key, value)
        
        if '{Tabela1}' in slide_replacements and tables:
            df = tables[0]
            rows, cols = df.shape
            table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(2.0), Inches(9.0), Inches(1.0)).table
            
            # Definir nomes das colunas
            for col in range(cols):
                table.cell(0, col).text = df.columns[col]
            
            # Preencher tabela com dados
            for row in range(1, rows+1):
                for col in range(cols):
                    table.cell(row, col).text = str(df.iat[row-1, col])
    
    presentation.save(output_path)

#LOCAL
#LOCAL DO WORD
source_doc_path = 'C:/Users/Gabriel/Desktop/Guia de Investimentos (lista de Trends).docx'
#LOCAL DO TEMPLATE
ppt_template_path = 'C:/Users/Gabriel/Desktop/template.pptx'
#LOCAL DE SAIDA 
output_ppt_path = 'C:/Users/Gabriel/Desktop/output_presentation.pptx'
#local de saida
output_pdf_path = 'C:/Users/Gabriel/Desktop/output_presentation.pdf'

#PARAGRAFOS
paragraphs, tables = extract_doc_contents(source_doc_path)

#Texto Tabelas
table_text = ""
if tables:
    table_text = "\n".join(["\t".join(row) for row in tables[0].values.astype(str)])

#Dia de Hj

current_date = datetime.now().strftime('%d/%m/%Y')


# Layouts Possivies
def get_slide_replacements(paragraphs, start, end):
    return {
        '{Titulo}': paragraphs[start] if len(paragraphs) > start else '',
        '{Subtitulo}': paragraphs[start + 1] if len(paragraphs) > start + 1 else '',
        '{Corpo}': "\n\n".join(paragraphs[start + 2:end]) if len(paragraphs) > start + 2 else '',
        '{date}': current_date,
    }


# Definir substituições para cada slide
#MUDAR SE NECESSARIO
replacements = [
    get_slide_replacements(paragraphs, 0, 16),
    get_slide_replacements(paragraphs, 16, 20),
    get_slide_replacements(paragraphs, 20, 28),
    get_slide_replacements(paragraphs, 28, 38),
    get_slide_replacements(paragraphs, 38, 47),
    get_slide_replacements(paragraphs, 47, 56),
    get_slide_replacements(paragraphs, 56, 64),
    get_slide_replacements(paragraphs, 64, 76),
    get_slide_replacements(paragraphs, 76, 86),
    get_slide_replacements(paragraphs, 86, 95),
]

# Substituir placeholders no Template
substitute_placeholders_ppt(ppt_template_path, output_ppt_path, replacements, tables)
